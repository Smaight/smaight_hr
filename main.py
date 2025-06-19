import os
import subprocess
import PyPDF2
import olefile
import zlib
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
import argparse
import openpyxl
import xlrd
from pptx import Presentation


def extract_text_from_ppt(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = None
        print(f"ppt 변환 오류: {e}")
    return text


def convert_ppt_to_pptx_with_libreoffice(ppt_path):
    output_dir = os.path.dirname(ppt_path)
    result = subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pptx",
            "--outdir",
            output_dir,
            ppt_path,
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode == 0:
        base = os.path.splitext(os.path.basename(ppt_path))[0]
        return os.path.join(output_dir, base + ".pptx")
    else:
        print(f"libreoffice ppt→pptx 변환 실패: {result.stderr}")
        return None


def extract_text_from_excel(file_path):
    text = ""
    ext = file_path.lower().split(".")[-1]
    try:
        if ext == "xlsx":
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        [str(cell) if cell is not None else "" for cell in row]
                    )
                    text += row_text + "\n"
        elif ext == "xls":
            wb = xlrd.open_workbook(file_path)
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_text = "\t".join(
                        [str(cell) if cell is not None else "" for cell in row]
                    )
                    text += row_text + "\n"
        else:
            return None
    except Exception as e:
        print(f"엑셀 변환 오류: {e}")
        return None
    return text


def convert_doc_to_docx_with_libreoffice(doc_path):
    """
    리눅스에서 libreoffice를 이용해 .doc 파일을 .docx로 변환
    """
    try:
        output_dir = os.path.dirname(doc_path)
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                output_dir,
                doc_path,
            ],
            capture_output=True,
            text=True,
        )
        if result.returncode == 0:
            base = os.path.splitext(os.path.basename(doc_path))[0]
            return os.path.join(output_dir, base + ".docx")
        else:
            print(f"libreoffice 변환 실패: {result.stderr}")
            return None
    except Exception as e:
        print(f"libreoffice 변환 예외: {e}")
        return None


def iter_block_items(parent):
    """
    문서의 block-level 요소를 순서대로 yield함.
    표 또는 문단을 반환.
    """
    for child in parent.element.body.iterchildren():
        if isinstance(child, CT_P):  # 일반 단락
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):  # 표
            yield Table(child, parent)


def convert_files_to_text(source_folder, target_folder):
    if not os.path.exists(target_folder):
        os.makedirs(target_folder)
        print(f"대상 폴더 생성: {target_folder}")

    for filename in os.listdir(source_folder):
        source_filepath = os.path.join(source_folder, filename)

        if os.path.isfile(source_filepath):
            file_extension = filename.lower().split(".")[-1]
            text = ""
            converted = False

            try:
                print(f"변환 중... {filename}")

                if file_extension == "pdf":
                    with open(source_filepath, "rb") as pdf_file:
                        reader = PyPDF2.PdfReader(pdf_file)
                        for page in reader.pages:
                            extracted = page.extract_text()
                            text += extracted or ""
                    converted = True

                elif file_extension == "docx":
                    document = Document(source_filepath)

                    for block in iter_block_items(document):
                        if isinstance(block, Paragraph):
                            if block.text.strip():
                                text += block.text + "\n"
                        elif isinstance(block, Table):
                            for row in block.rows:
                                row_text = [cell.text.strip() for cell in row.cells]
                                text += "\t".join(row_text) + "\n"
                    converted = True

                elif file_extension == "doc":
                    print(f".doc → .docx 변환 시도: {filename}")
                    docx_path = convert_doc_to_docx_with_libreoffice(source_filepath)

                    if docx_path and os.path.exists(docx_path):
                        document = Document(docx_path)

                        for block in iter_block_items(document):
                            if isinstance(block, Paragraph):
                                if block.text.strip():
                                    text += block.text + "\n"
                            elif isinstance(block, Table):
                                for row in block.rows:
                                    row_text = [cell.text.strip() for cell in row.cells]
                                    text += "\t".join(row_text) + "\n"
                        converted = True
                        os.remove(docx_path)  # 변환된 .docx 임시파일 삭제
                    else:
                        print(f"{filename} 변환 실패")

                elif file_extension in ["xlsx", "xls"]:
                    text = extract_text_from_excel(source_filepath)
                    if text is not None:
                        converted = True

                elif file_extension == "ppt":
                    pptx_path = convert_ppt_to_pptx_with_libreoffice(source_filepath)
                    if pptx_path and os.path.exists(pptx_path):
                        text = extract_text_from_ppt(pptx_path)
                        converted = True
                        os.remove(pptx_path)
                    else:
                        print(f"{filename} 변환 실패")

                elif file_extension in ["pptx"]:
                    text = extract_text_from_ppt(source_filepath)
                    if text is not None:
                        converted = True

                else:
                    print(f"지원 안 함: {filename}")
                    continue

                if converted:
                    base_filename, _ = os.path.splitext(filename)
                    target_filepath = os.path.join(
                        target_folder, f"{base_filename}.txt"
                    )
                    with open(target_filepath, "w", encoding="utf-8-sig") as txt_file:
                        txt_file.write(text)
                    print(f"변환 완료 → {base_filename}.txt")

            except Exception as e:
                print(f"{filename} 변환 오류: {e}")

    print("\n모든 파일 변환이 완료되었습니다!")


# 기존 함수와 import는 그대로 두고, 아래만 추가/수정
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert files to text")
    parser.add_argument(
        "--source", type=str, required=True, help="Source directory path"
    )
    parser.add_argument(
        "--target", type=str, required=True, help="Target directory path"
    )
    args = parser.parse_args()

    source_directory = args.source
    target_directory = args.target

    convert_files_to_text(source_directory, target_directory)
